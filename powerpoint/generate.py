"""
StockRAG Progress Presentation Generator — Update 2
Black & Green theme | CSE 5914 Capstone
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Colors ──────────────────────────────────────────────────────────────────
BLACK       = RGBColor(0x0A, 0x0A, 0x0A)
DARK_CARD   = RGBColor(0x14, 0x14, 0x14)
GREEN       = RGBColor(0x00, 0xC8, 0x53)
GREEN_DIM   = RGBColor(0x00, 0x7A, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_GRAY  = RGBColor(0x77, 0x77, 0x77)
DONE_GREEN  = RGBColor(0x00, 0xE6, 0x76)
TODO_AMBER  = RGBColor(0xFF, 0xD7, 0x40)
CYAN        = RGBColor(0x00, 0xD4, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BLACK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def accent_bar(slide, y_pos=Inches(0.62), height=Inches(0.04)):
    add_rect(slide, 0, y_pos, SLIDE_W, height, GREEN)


def slide_header(slide, title, subtitle=None):
    accent_bar(slide)
    add_text(slide, title,
             left=Inches(0.5), top=Inches(0.08),
             width=Inches(9), height=Inches(0.55),
             font_size=28, bold=True, color=GREEN)
    if subtitle:
        add_text(slide, subtitle,
                 left=Inches(0.5), top=Inches(0.68),
                 width=Inches(12), height=Inches(0.35),
                 font_size=13, color=GRAY, italic=True)


def tag(slide, label, color=GREEN_DIM, text_color=GREEN,
        left=Inches(10.3), top=Inches(0.1)):
    add_rect(slide, left, top, Inches(2.7), Inches(0.4), color)
    add_text(slide, label,
             left=left, top=top,
             width=Inches(2.7), height=Inches(0.4),
             font_size=11, bold=True, color=text_color,
             align=PP_ALIGN.CENTER)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def done_todo_cols(slide, done_items, todo_items,
                   done_label="COMPLETED", todo_label="IN PROGRESS"):
    col_top  = Inches(1.15)
    col_h    = Inches(5.7)
    lbl_h    = Inches(0.4)
    body_top = col_top + lbl_h + Inches(0.08)
    body_h   = col_h - lbl_h - Inches(0.08)

    # Done column
    add_rect(slide, Inches(0.35), col_top, Inches(5.9), col_h, DARK_CARD)
    add_rect(slide, Inches(0.35), col_top, Inches(5.9), Inches(0.04), DONE_GREEN)
    add_text(slide, done_label,
             Inches(0.55), col_top + Inches(0.08),
             Inches(5.5), lbl_h,
             font_size=12, bold=True, color=DONE_GREEN)
    tb = slide.shapes.add_textbox(Inches(0.55), body_top, Inches(5.5), body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in done_items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE

    # Todo column
    add_rect(slide, Inches(6.6), col_top, Inches(6.4), col_h, DARK_CARD)
    add_rect(slide, Inches(6.6), col_top, Inches(6.4), Inches(0.04), TODO_AMBER)
    add_text(slide, todo_label,
             Inches(6.8), col_top + Inches(0.08),
             Inches(6.0), lbl_h,
             font_size=12, bold=True, color=TODO_AMBER)
    tb2 = slide.shapes.add_textbox(Inches(6.8), body_top, Inches(6.0), body_h)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for item in todo_items:
        if first:
            p = tf2.paragraphs[0]; first = False
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Green accent block (right side)
    add_rect(s, Inches(9.8), 0, Inches(3.53), Inches(7.5), RGBColor(0x00, 0x3A, 0x18))

    # Thin left border bar
    add_rect(s, 0, 0, Inches(0.08), SLIDE_H, GREEN)

    # Bottom accent line
    add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.04), GREEN_DIM)

    # App name
    add_text(s, "StockRAG", Inches(0.5), Inches(1.4),
             Inches(9), Inches(1.6),
             font_size=80, bold=True, color=GREEN)

    # Tagline
    add_text(s, "AI-Powered Stock Discovery",
             Inches(0.5), Inches(2.95), Inches(9), Inches(0.7),
             font_size=28, color=WHITE)

    # Sub-line
    add_text(s, "CSE 5914 Capstone  ·  Progress Update 2  ·  Spring 2026",
             Inches(0.5), Inches(3.6), Inches(9.2), Inches(0.45),
             font_size=15, color=GRAY)

    # Divider line
    add_rect(s, Inches(0.5), Inches(4.1), Inches(5), Inches(0.03), GREEN_DIM)

    # Team line
    add_text(s, "Michael Gorman  ·  Alex Fizer  ·  Arnav Rajashekara  ·  Manu  ·  Krish Patel  ·  Jay Jivandas",
             Inches(0.5), Inches(4.25), Inches(9.1), Inches(0.4),
             font_size=12, color=GRAY)

    # Right-column label
    add_text(s, "PROGRESS\nUPDATE", Inches(10.1), Inches(2.8),
             Inches(3), Inches(1.5),
             font_size=32, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Welcome everyone. This is our second progress update for StockRAG, "
        "our AI-powered stock discovery tool. Today we'll walk through what each "
        "team has accomplished since the last update and show a live demo of the frontend.")


def slide_02_recap(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Quick Recap",
                 "What StockRAG does in 30 seconds")

    # Three pillars with updated descriptions
    cols = [
        ("FRONTEND",
         "React chat UI where users\nask natural-language\nstock questions",
         Inches(0.5), RGBColor(0x00, 0x3A, 0x18)),
        ("BACKEND",
         "FastAPI server routing\nqueries to the RAG engine\nand returning results",
         Inches(4.6), RGBColor(0x00, 0x30, 0x15)),
        ("RAG PIPELINE",
         "SEC + Finnhub data fed into\nChromaDB embeddings,\nqueried via Gemini LLM",
         Inches(8.7), RGBColor(0x00, 0x26, 0x12)),
    ]
    for title, body, left, bg in cols:
        add_rect(s, left, Inches(1.15), Inches(3.9), Inches(2.3), bg)
        add_rect(s, left, Inches(1.15), Inches(3.9), Inches(0.04), GREEN)
        add_text(s, title, left + Inches(0.15), Inches(1.3),
                 Inches(3.6), Inches(0.5),
                 font_size=14, bold=True, color=GREEN)
        add_text(s, body, left + Inches(0.15), Inches(1.8),
                 Inches(3.6), Inches(1.2),
                 font_size=13, color=WHITE)

    # Flow summary
    add_rect(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.9), RGBColor(0x0F, 0x0F, 0x0F))
    add_rect(s, Inches(0.5), Inches(3.7), Inches(0.06), Inches(0.9), GREEN)
    add_text(s, 'User asks "EV stocks with high growth"  →  Backend queries RAG pipeline  →  '
                'ChromaDB similarity search  →  Gemini generates explanation  →  StockCards returned',
             Inches(0.75), Inches(3.78), Inches(11.9), Inches(0.75),
             font_size=13, color=WHITE)

    # Tech stack chips
    stack = ["React + TypeScript", "Vite", "Mantine UI", "FastAPI", "Pydantic",
             "ChromaDB", "all-MiniLM-L6-v2", "Google Gemini", "Finnhub API"]
    x = Inches(0.5)
    for label in stack:
        w = Inches(len(label) * 0.11 + 0.3)
        add_rect(s, x, Inches(4.85), w, Inches(0.35), RGBColor(0x1A, 0x1A, 0x1A))
        add_text(s, label, x, Inches(4.85), w, Inches(0.35),
                 font_size=9, color=GREEN_DIM, align=PP_ALIGN.CENTER)
        x += w + Inches(0.12)

    add_speaker_notes(s,
        "Quick refresher for anyone who missed the first update. StockRAG has three main pieces: "
        "a React chat frontend, a FastAPI backend, and a RAG pipeline built on SEC EDGAR and Finnhub data. "
        "The user types a natural language query, it flows through the pipeline, and they get back "
        "AI-generated stock recommendations with explanations.")


def slide_03_frontend(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Frontend Progress", "Michael Gorman & Alex Fizer")
    tag(s, "React + TypeScript + Mantine")

    done = [
        "Complete chat interface — AppHeader, MainLayout, ChatInput",
        "WelcomeScreen with example query chips",
        "MessageList with smooth auto-scroll behavior",
        "StockCard component: ticker, match %, sector, metrics, 'Why it fits'",
        "Loading skeletons and fade-in animations",
        "Custom dark/green Mantine theme system",
        "useChat hook with optimistic UI and error states",
        "useAutoScroll and useChatInput custom hooks",
        "Axios API client with automatic mock fallback",
        "TypeScript interfaces aligned with backend Pydantic models",
        "Enforced code rules — <15 line functions, no prop drilling",
    ]
    todo = [
        "Connect to live backend API (remove mocks)",
        "Login & Signup pages",
        "User profile page with settings",
        "Portfolio state management (Zustand/Context)",
        "Portfolio dashboard with charts (Recharts)",
        "Error toast notification system",
    ]
    done_todo_cols(s, done, todo)

    add_speaker_notes(s,
        "The frontend is the most mature piece right now. We have a fully working chat interface "
        "with React, TypeScript, and Mantine UI. Users can type queries, see loading skeletons, "
        "and get back beautifully rendered StockCards showing ticker, match percentage, sector, "
        "and a 'Why it fits' explanation. Right now it runs on mock data — the main next step "
        "is connecting to the live backend once it's ready. After that we'll add auth pages "
        "and a portfolio dashboard. We'll demo this live in a moment.")


def slide_04_backend(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Backend Progress", "Arnav Rajashekara & Manu")
    tag(s, "FastAPI + Pydantic + ChromaDB")

    done = [
        "FastAPI app scaffolded — main.py, config.py, route structure",
        "Pydantic v2 settings: Gemini key, ChromaDB dir, port config",
        "Full API models: RecommendationRequest, StockRecommendation,",
        "    RecommendationResponse, HealthCheckResponse, StatsResponse",
        "Complete route implementation: GET /health, GET /stats,",
        "    POST /api/recommendations with full request handling",
        "Service layer stubs: llm_service.py, vector_db.py",
        "Frontend-aligned API contracts (TypeScript ↔ Pydantic)",
        "Jupyter notebook for model validation testing",
        "Docker Compose configuration",
    ]
    todo = [
        "Implement LLM service (Gemini primary, Groq fallback)",
        "Wire ChromaDB vector DB into service layer",
        "Complete RAG pipeline orchestration (pipeline.py)",
        "Embedding generation module (embeddings.py)",
        "Retrieval + similarity search logic (retrieval.py)",
        "Data enrichment via yfinance + Finnhub",
        "End-to-end integration test (mock → live)",
    ]
    done_todo_cols(s, done, todo)

    add_speaker_notes(s,
        "The backend team has the full FastAPI structure in place. All the route definitions, "
        "data models, and API contracts are done and aligned with what the frontend expects. "
        "The routes.py file already imports from the RAG pipeline module — it just needs the "
        "actual implementation wired in. The big remaining work is implementing the LLM service "
        "to talk to Gemini, connecting ChromaDB, and building out the retrieval logic. "
        "Once the RAG team has data in ChromaDB, this should come together quickly.")


def slide_05_rag(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "RAG Pipeline Progress", "Krish Patel & Jay Jivandas")
    tag(s, "SEC EDGAR + Finnhub + Gemini")

    done = [
        "Analyzed 19,215 SEC EDGAR company filings (18.4 GB raw data)",
        "Mapped key financial metrics: balance sheet, income, cash flow",
        "Documented taxonomies: us-gaap (~85%), dei, srt, invest, ifrs",
        "Flagged 2,647 empty files to skip in pipeline",
        "Built Finnhub data puller (546 lines) — pulls all free-tier data",
        "Tested Finnhub API for AAPL: profile, quotes, financials, news,",
        "    earnings, insider sentiment, SEC filings, recommendations",
        "Generated structured CSV + JSON output with auto-reporting",
        "Full backend folder structure scaffolded",
    ]
    todo = [
        "Parse + clean SEC EDGAR JSON at scale",
        "Generate embeddings using all-MiniLM-L6-v2",
        "Populate ChromaDB with company vectors",
        "Build similarity search / retrieval logic",
        "Integrate Finnhub for real-time price enrichment",
        "Prompt template engineering for Gemini",
        "End-to-end pipeline test: query → ranked results",
    ]
    done_todo_cols(s, done, todo)

    add_speaker_notes(s,
        "The RAG team has made solid progress since last time. The big new development is "
        "the Finnhub data puller — a 546-line script that pulls all free-tier endpoints for "
        "any given stock: company profiles, real-time quotes, financial statements, news, "
        "earnings data, insider sentiment, and more. They tested it against AAPL and documented "
        "exactly which endpoints work on the free tier vs which return 403s. "
        "The SEC EDGAR analysis from before is still the foundation — 19k+ companies worth of data. "
        "The critical next step is generating the embeddings and getting them into ChromaDB.")


def slide_06_integration(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Integration Status",
                 "Where the three workstreams connect")

    # Status bar for each connection point
    connections = [
        ("Frontend ↔ Backend",
         "READY TO CONNECT",
         "TypeScript interfaces match Pydantic models. Switching VITE_API_BASE_URL goes live.",
         Inches(1.15), DONE_GREEN),
        ("Backend ↔ RAG Pipeline",
         "CONTRACTS DEFINED",
         "routes.py imports rag.pipeline — service stubs exist, implementation next.",
         Inches(2.65), TODO_AMBER),
        ("RAG ↔ Data Sources",
         "DATA COLLECTED",
         "19k+ SEC EDGAR filings analyzed. Finnhub puller tested. Embeddings generation next.",
         Inches(4.15), TODO_AMBER),
    ]
    for label, status, desc, top, status_color in connections:
        add_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.2), DARK_CARD)
        add_rect(s, Inches(0.5), top, Inches(0.06), Inches(1.2), status_color)
        add_text(s, label, Inches(0.75), top + Inches(0.1),
                 Inches(4), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)
        add_text(s, status, Inches(5.5), top + Inches(0.12),
                 Inches(2.5), Inches(0.35),
                 font_size=10, bold=True, color=status_color)
        add_text(s, desc, Inches(0.75), top + Inches(0.55),
                 Inches(11.5), Inches(0.5),
                 font_size=12, color=GRAY)

    # Critical path
    add_rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2), RGBColor(0x0F, 0x0F, 0x0F))
    add_rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.04), GREEN)
    add_text(s, "CRITICAL PATH", Inches(0.7), Inches(5.8),
             Inches(3), Inches(0.35),
             font_size=11, bold=True, color=GREEN)

    # Path steps
    steps = ["SEC Data Parsing", "Embeddings", "ChromaDB", "Backend Services", "Live Frontend"]
    x = Inches(0.7)
    for i, step in enumerate(steps):
        sw = Inches(2.1)
        add_rect(s, x, Inches(6.2), sw, Inches(0.45), RGBColor(0x00, 0x3A, 0x18))
        add_text(s, step, x, Inches(6.2), sw, Inches(0.45),
                 font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        x += sw + Inches(0.05)
        if i < len(steps) - 1:
            add_rect(s, x, Inches(6.38), Inches(0.2), Inches(0.06), GREEN)
            x += Inches(0.25)

    add_speaker_notes(s,
        "Here's where we stand on integration. The frontend-to-backend contract is fully aligned — "
        "TypeScript types match Pydantic models exactly, so flipping one env variable connects them. "
        "The backend has route stubs that import the RAG pipeline module, so once the pipeline is "
        "implemented it plugs right in. The critical path right now goes through the RAG team: "
        "parse the SEC data, generate embeddings, load ChromaDB, then the backend team wires it up, "
        "and the frontend goes live. Each step unblocks the next.")


def slide_07_live_demo(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Large green border frame
    add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x0F, 0x0F, 0x0F))
    add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(0.05), GREEN)
    add_rect(s, Inches(0.3), Inches(0.3), Inches(0.05), Inches(6.9), GREEN)
    add_rect(s, Inches(0.3), Inches(7.15), Inches(12.73), Inches(0.05), GREEN)
    add_rect(s, Inches(12.98), Inches(0.3), Inches(0.05), Inches(6.9), GREEN)

    # Center content
    add_text(s, "LIVE DEMO", Inches(0.5), Inches(2.0),
             Inches(12.3), Inches(1.2),
             font_size=64, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)

    add_text(s, "Frontend Chat Interface",
             Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.6),
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(s, Inches(5), Inches(4.1), Inches(3.3), Inches(0.03), GREEN_DIM)

    # What we'll show
    demo_items = [
        "Type a natural-language stock query",
        "Watch loading skeletons + animations",
        "See StockCard results with AI explanations",
    ]
    y = Inches(4.4)
    for item in demo_items:
        add_text(s, "▸  " + item,
                 Inches(3.5), y, Inches(6.3), Inches(0.4),
                 font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
        y += Inches(0.4)

    # Subtle corner accent
    add_rect(s, Inches(11.5), Inches(5.8), Inches(1.3), Inches(1.0), RGBColor(0x00, 0x3A, 0x18))
    add_text(s, "EXIT\nSLIDES", Inches(11.5), Inches(5.9),
             Inches(1.3), Inches(0.8),
             font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "At this point we'll exit the slideshow and switch to a live demo of the frontend. "
        "We'll show the chat interface running locally with Vite, type in a stock query like "
        "'EV companies with strong growth', and walk through the full user experience: "
        "the welcome screen with example chips, the loading skeleton animations, "
        "and the StockCard results with ticker, match percentage, sector info, and the "
        "'Why it fits' explanation. Right now it's running on mock data but the UI behavior "
        "is identical to what it will look like with real backend responses.")


def slide_08_next_steps(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "What's Next", "Priorities for each team before the next check-in")

    teams = [
        ("Frontend",
         ["Wire up live backend API",
          "Build login + signup pages",
          "Add user profile page",
          "Portfolio dashboard (Recharts)",
          "Error toast notifications"],
         Inches(0.35), RGBColor(0x00, 0x3A, 0x18)),
        ("Backend",
         ["Implement Gemini LLM service",
          "Connect ChromaDB vector store",
          "Wire RAG pipeline into routes",
          "Add Finnhub data enrichment",
          "End-to-end integration test"],
         Inches(4.55), RGBColor(0x00, 0x30, 0x15)),
        ("RAG Pipeline",
         ["Parse SEC EDGAR JSON at scale",
          "Generate MiniLM embeddings",
          "Load vectors into ChromaDB",
          "Build retrieval / search logic",
          "Gemini prompt engineering"],
         Inches(8.75), RGBColor(0x00, 0x26, 0x12)),
    ]
    for title, items, left, bg in teams:
        add_rect(s, left, Inches(1.15), Inches(3.9), Inches(4.5), DARK_CARD)
        add_rect(s, left, Inches(1.15), Inches(3.9), Inches(0.5), bg)
        add_rect(s, left, Inches(1.15), Inches(3.9), Inches(0.04), GREEN)
        add_text(s, title, left + Inches(0.15), Inches(1.2),
                 Inches(3.6), Inches(0.45),
                 font_size=15, bold=True, color=WHITE)
        tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.8),
                                   Inches(3.6), Inches(3.7))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(10)
            run = p.add_run()
            run.text = "→  " + item
            run.font.size = Pt(13)
            run.font.color.rgb = WHITE

    # Bottom note
    add_rect(s, Inches(0.35), Inches(5.9), Inches(12.3), Inches(0.7), RGBColor(0x0F, 0x0F, 0x0F))
    add_rect(s, Inches(0.35), Inches(5.9), Inches(0.06), Inches(0.7), CYAN)
    add_text(s, "Goal:  RAG team unblocks backend → backend unblocks frontend → full pipeline demo at next update",
             Inches(0.6), Inches(5.98), Inches(11.8), Inches(0.45),
             font_size=13, bold=True, color=CYAN)

    add_speaker_notes(s,
        "Here's the game plan going forward. The RAG team is the critical path — once they "
        "get SEC data parsed and embeddings into ChromaDB, the backend team can wire up the "
        "LLM service and retrieval logic. Once the backend is returning real results, the "
        "frontend team flips the API URL and we have a working end-to-end pipeline. "
        "In parallel, the frontend team will build out auth pages and the portfolio dashboard. "
        "The goal for the next update is a fully working demo with real data flowing through.")


def slide_09_closing(prs):
    s = blank_slide(prs)
    fill_bg(s)

    add_rect(s, 0, 0, Inches(0.08), SLIDE_H, GREEN)
    add_rect(s, Inches(9.8), 0, Inches(3.53), SLIDE_H, RGBColor(0x00, 0x3A, 0x18))

    # Bottom accent
    add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.04), GREEN_DIM)

    add_text(s, "Questions?",
             Inches(0.5), Inches(2.2), Inches(9), Inches(1.4),
             font_size=72, bold=True, color=GREEN)

    add_text(s, "StockRAG  ·  CSE 5914 Capstone  ·  Spring 2026",
             Inches(0.5), Inches(3.6), Inches(9), Inches(0.45),
             font_size=16, color=GRAY)

    add_rect(s, Inches(0.5), Inches(4.15), Inches(4), Inches(0.03), GREEN_DIM)

    add_text(s, "github.com/jjivandas/CSE-5914-Capstone",
             Inches(0.5), Inches(4.25), Inches(9), Inches(0.4),
             font_size=12, color=GRAY, italic=True)

    add_text(s, "THANK\nYOU", Inches(10.1), Inches(2.8),
             Inches(3), Inches(1.5),
             font_size=32, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "That's our update. We're happy to take any questions about the architecture, "
        "the tech stack, or our timeline. The code is all on GitHub if you want to dig in.")


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_01_title(prs)        # 1. Title
    slide_02_recap(prs)        # 2. Quick recap
    slide_03_frontend(prs)     # 3. Frontend progress
    slide_04_backend(prs)      # 4. Backend progress
    slide_05_rag(prs)          # 5. RAG progress
    slide_06_integration(prs)  # 6. Integration status
    slide_07_live_demo(prs)    # 7. Live demo (exit slides here)
    slide_08_next_steps(prs)   # 8. Next steps
    slide_09_closing(prs)      # 9. Questions

    out = "stockrag_progress_2.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
