"""Generate RAG Pipeline Progress slide as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── colours ───────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0d, 0x0d, 0x1a)   # near-black blue
CARD_BG     = RGBColor(0x12, 0x12, 0x2e)   # slightly lighter card
GREEN       = RGBColor(0x00, 0xff, 0x88)   # bright green (title, badge, done)
YELLOW      = RGBColor(0xff, 0xcc, 0x00)   # yellow (in-progress)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT_GREY  = RGBColor(0xcc, 0xcc, 0xcc)
DARK_GREEN  = RGBColor(0x00, 0x44, 0x22)   # badge background


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, left, top, width, height, items, color=LIGHT_GREY, font_size=11):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"\u25ba  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def build_slide():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    W = prs.slide_width
    H = prs.slide_height

    # ── background ────────────────────────────────────────────────────────────
    bg = add_rect(slide, 0, 0, W, H, BG)

    # ── title ─────────────────────────────────────────────────────────────────
    add_textbox(slide,
                Inches(0.4), Inches(0.25),
                Inches(8), Inches(0.7),
                "RAG Pipeline Progress",
                font_size=32, bold=True, color=GREEN)

    # ── authors ───────────────────────────────────────────────────────────────
    add_textbox(slide,
                Inches(0.4), Inches(0.95),
                Inches(5), Inches(0.35),
                "Krish Patel & Jay Jivandas",
                font_size=11, color=LIGHT_GREY)

    # ── badge: SEC EDGAR + Finnhub + Gemini ───────────────────────────────────
    badge_w, badge_h = Inches(2.6), Inches(0.38)
    badge_l = W - badge_w - Inches(0.3)
    badge_t = Inches(0.28)
    add_rect(slide, badge_l, badge_t, badge_w, badge_h,
             DARK_GREEN, line_color=GREEN, line_width=Pt(1.5))
    add_textbox(slide, badge_l, badge_t + Inches(0.04),
                badge_w, badge_h,
                "SEC EDGAR + Finnhub + Gemini",
                font_size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ── divider line ──────────────────────────────────────────────────────────
    line = slide.shapes.add_shape(1,
        Inches(0.4), Inches(1.35),
        Inches(4.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()

    # ── COMPLETED card ────────────────────────────────────────────────────────
    card_t = Inches(1.5)
    card_h = Inches(5.7)
    card_w = Inches(5.8)

    add_rect(slide, Inches(0.3), card_t, card_w, card_h,
             CARD_BG, line_color=GREEN, line_width=Pt(1.2))

    add_textbox(slide,
                Inches(0.5), card_t + Inches(0.15),
                Inches(3), Inches(0.4),
                "COMPLETED",
                font_size=13, bold=True, color=GREEN)

    completed_items = [
        "Analyzed 19,215 SEC EDGAR company filings (18.4 GB raw data)",
        "Mapped key financial metrics: balance sheet, income, cash flow",
        "Documented taxonomies: us-gaap (~85%), dei, srt, invest, ifrs",
        "Flagged 2,647 empty files to skip in pipeline",
        "Built Finnhub data puller — pulls all free-tier data",
        "Parsed + cleaned EDGAR JSON → 4 normalized Parquet tables (521 MB)",
        "Generated per-company RAG JSON files for all 7,138 companies (6.6 GB)",
        "Embedded companies using all-MiniLM-L6-v2 into ChromaDB\n"
        "   (7,138 profiles + 28,124 annual snapshots)",
        "Built two-stage RAG pipeline: query expansion → retrieval (top 20)\n"
        "   → LLM reranking (top 10) → generation",
        "Integrated Gemini (gemini-2.5-flash-lite) for expansion,\n"
        "   reranking, and generation",
        "Enriched company profiles with Wikipedia/Finnhub descriptions",
        "End-to-end pipeline verified: query → ranked results + citations",
        "Built demo UI for live presentation",
    ]

    add_bullet_box(slide,
                   Inches(0.5), card_t + Inches(0.6),
                   card_w - Inches(0.4), card_h - Inches(0.75),
                   completed_items, color=LIGHT_GREY, font_size=10)

    # ── IN PROGRESS card ──────────────────────────────────────────────────────
    card2_l = Inches(6.5)
    card2_w = Inches(6.5)
    card2_h = Inches(5.7)

    add_rect(slide, card2_l, card_t, card2_w, card2_h,
             CARD_BG, line_color=YELLOW, line_width=Pt(1.2))

    add_textbox(slide,
                card2_l + Inches(0.2), card_t + Inches(0.15),
                Inches(3), Inches(0.4),
                "IN PROGRESS",
                font_size=13, bold=True, color=YELLOW)

    inprogress_items = [
        "Re-index ChromaDB with company descriptions\n"
        "   (Wikipedia/Finnhub enrichment — script done, re-run pending)",
        "Integrate real-time stock price data\n"
        "   (Finnhub live prices not yet connected to pipeline)",
        "React frontend\n"
        "   (component stubs exist in stockrag/frontend/, not yet implemented)",
        "RAG evaluation harness\n"
        "   (hit rate, MRR, faithfulness scoring not yet built)",
    ]

    add_bullet_box(slide,
                   card2_l + Inches(0.2), card_t + Inches(0.6),
                   card2_w - Inches(0.4), card2_h - Inches(0.75),
                   inprogress_items, color=LIGHT_GREY, font_size=11)

    # ── save ──────────────────────────────────────────────────────────────────
    out = "/home/krish/uni/cse5914/CSE-5914-Capstone/stockrag/demo/rag_pipeline_progress.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build_slide()
